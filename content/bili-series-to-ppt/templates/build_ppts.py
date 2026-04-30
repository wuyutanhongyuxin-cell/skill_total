"""Render each chapter outline JSON in outlines/ to a .pptx in ppt/.

Run:  uv run --project bili2text python scripts/build_ppts.py

Reads:  outlines/p*.json  (validated upstream by fix_outline_quotes.py)
Writes: ppt/p<NN>.pptx

Slide structure per chapter (strictly fixed):
  1. Title slide           : "第NN讲  <title>" + subtitle + footer
  2. Learning objectives   : objectives bullets
  3. Outline               : outline bullets
  4. For each section in sections[]:
       a) Bullets slide    : section.bullets
       b) Examples slide   : section.examples (only if non-empty)
       c) Tips slide       : section.tips     (only if non-empty)
  5. Exercises             : exercises[] (only if non-empty)
  6. Summary               : summary[]   (only if non-empty)

Layout: 16:9 widescreen (13.33 x 7.5 in).
Colors: navy / accent / gray / light (canonical palette — do not extend).
Fonts:  unified font sizes per slide type (title 28pt / body 18pt etc.).
"""
from __future__ import annotations
import json, pathlib
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

ROOT = pathlib.Path(__file__).resolve().parents[1]
OUTLINES = ROOT / "outlines"
PPT_DIR = ROOT / "ppt"
PPT_DIR.mkdir(exist_ok=True)

NAVY = RGBColor(0x14, 0x2A, 0x5C)
ACCENT = RGBColor(0xB8, 0x3A, 0x3A)
GRAY = RGBColor(0x55, 0x55, 0x55)
LIGHT = RGBColor(0xF2, 0xEE, 0xE3)


def add_title_slide(prs, p, full_title, subtitle):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg = s.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = LIGHT
    tb = s.shapes.add_textbox(Inches(0.6), Inches(2.0), Inches(12.0), Inches(1.2))
    tf = tb.text_frame
    tf.word_wrap = True
    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.LEFT
    r = p1.add_run()
    r.text = f"第{p:02d}讲  {full_title}"
    r.font.size = Pt(36)
    r.font.bold = True
    r.font.color.rgb = NAVY
    sb = s.shapes.add_textbox(Inches(0.6), Inches(3.4), Inches(12.0), Inches(0.8))
    sbtf = sb.text_frame
    sbtf.word_wrap = True
    sp = sbtf.paragraphs[0]
    sr = sp.add_run()
    sr.text = subtitle or "教学课件"
    sr.font.size = Pt(20)
    sr.font.color.rgb = GRAY
    fb = s.shapes.add_textbox(Inches(0.6), Inches(6.6), Inches(12.0), Inches(0.4))
    ft = fb.text_frame
    fp = ft.paragraphs[0]
    fr = fp.add_run()
    fr.text = "系列课程 / 教学讲义"
    fr.font.size = Pt(12)
    fr.font.color.rgb = GRAY


def add_bullet_slide(prs, heading, bullets, eyebrow=None):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    if eyebrow:
        eb = s.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12.5), Inches(0.4))
        et = eb.text_frame
        ep = et.paragraphs[0]
        er = ep.add_run()
        er.text = eyebrow
        er.font.size = Pt(12)
        er.font.color.rgb = ACCENT
        er.font.bold = True
    hb = s.shapes.add_textbox(Inches(0.5), Inches(0.55), Inches(12.5), Inches(0.9))
    ht = hb.text_frame
    ht.word_wrap = True
    hp = ht.paragraphs[0]
    hr = hp.add_run()
    hr.text = heading
    hr.font.size = Pt(28)
    hr.font.bold = True
    hr.font.color.rgb = NAVY
    bb = s.shapes.add_textbox(Inches(0.6), Inches(1.55), Inches(12.4), Inches(5.6))
    bt = bb.text_frame
    bt.word_wrap = True
    for i, b in enumerate(bullets):
        para = bt.paragraphs[0] if i == 0 else bt.add_paragraph()
        para.alignment = PP_ALIGN.LEFT
        if isinstance(b, dict):
            label = b.get("label") or ""
            text = b.get("text") or ""
            if label:
                r = para.add_run()
                r.text = f"  ▸ {label}  "
                r.font.size = Pt(18)
                r.font.bold = True
                r.font.color.rgb = ACCENT
            r2 = para.add_run()
            r2.text = text
            r2.font.size = Pt(18)
            r2.font.color.rgb = NAVY
        else:
            r = para.add_run()
            r.text = f"  •  {b}"
            r.font.size = Pt(18)
            r.font.color.rgb = NAVY
        para.space_after = Pt(8)


def add_example_slide(prs, heading, examples, eyebrow="例题 / 例文"):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    eb = s.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12.5), Inches(0.4))
    et = eb.text_frame
    ep = et.paragraphs[0]
    er = ep.add_run()
    er.text = eyebrow
    er.font.size = Pt(12)
    er.font.color.rgb = ACCENT
    er.font.bold = True
    hb = s.shapes.add_textbox(Inches(0.5), Inches(0.55), Inches(12.5), Inches(0.9))
    ht = hb.text_frame
    ht.word_wrap = True
    hp = ht.paragraphs[0]
    hr = hp.add_run()
    hr.text = heading
    hr.font.size = Pt(28)
    hr.font.bold = True
    hr.font.color.rgb = NAVY
    bb = s.shapes.add_textbox(Inches(0.6), Inches(1.55), Inches(12.4), Inches(5.6))
    bt = bb.text_frame
    bt.word_wrap = True
    for i, ex in enumerate(examples):
        jp = ex.get("jp", "")
        read = ex.get("read", "")
        note = ex.get("note", "")
        para = bt.paragraphs[0] if i == 0 else bt.add_paragraph()
        r1 = para.add_run()
        r1.text = f"  例 {i+1}.  "
        r1.font.size = Pt(18)
        r1.font.bold = True
        r1.font.color.rgb = ACCENT
        r2 = para.add_run()
        r2.text = jp
        r2.font.size = Pt(20)
        r2.font.bold = True
        r2.font.color.rgb = NAVY
        if read:
            r3 = para.add_run()
            r3.text = f"   〔読〕{read}"
            r3.font.size = Pt(16)
            r3.font.color.rgb = GRAY
        if note:
            sub = bt.add_paragraph()
            rs = sub.add_run()
            rs.text = f"        ↳ {note}"
            rs.font.size = Pt(15)
            rs.font.color.rgb = GRAY
        para.space_after = Pt(4)


def add_exercise_slide(prs, exercises, eyebrow="课堂练习"):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    eb = s.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12.5), Inches(0.4))
    er = eb.text_frame.paragraphs[0].add_run()
    er.text = eyebrow
    er.font.size = Pt(12)
    er.font.bold = True
    er.font.color.rgb = ACCENT
    hb = s.shapes.add_textbox(Inches(0.5), Inches(0.55), Inches(12.5), Inches(0.9))
    hp = hb.text_frame.paragraphs[0]
    hr = hp.add_run()
    hr.text = "练习题"
    hr.font.size = Pt(28)
    hr.font.bold = True
    hr.font.color.rgb = NAVY
    bb = s.shapes.add_textbox(Inches(0.6), Inches(1.55), Inches(12.4), Inches(5.6))
    bt = bb.text_frame
    bt.word_wrap = True
    for i, ex in enumerate(exercises):
        q = ex.get("q", "")
        a = ex.get("a", "")
        para = bt.paragraphs[0] if i == 0 else bt.add_paragraph()
        r = para.add_run()
        r.text = f"  Q{i+1}.  {q}"
        r.font.size = Pt(17)
        r.font.color.rgb = NAVY
        if a:
            sub = bt.add_paragraph()
            rs = sub.add_run()
            rs.text = f"        答：{a}"
            rs.font.size = Pt(15)
            rs.font.color.rgb = ACCENT
        para.space_after = Pt(6)


def add_summary_slide(prs, summary):
    add_bullet_slide(prs, "本讲小结", summary, eyebrow="Summary")


def build_one(outline_path: pathlib.Path) -> pathlib.Path:
    data = json.loads(outline_path.read_text(encoding="utf-8"))
    p = data["p"]
    title = data["title"]
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    add_title_slide(prs, p, title, data.get("subtitle", ""))
    if data.get("objectives"):
        add_bullet_slide(prs, "学习目标", data["objectives"], eyebrow="Learning Objectives")
    if data.get("outline"):
        add_bullet_slide(prs, "本讲提纲", data["outline"], eyebrow="Outline")
    for sec in data.get("sections", []):
        add_bullet_slide(prs, sec["heading"], sec.get("bullets", []), eyebrow=sec.get("eyebrow", "知识点"))
        if sec.get("examples"):
            add_example_slide(prs, f"{sec['heading']} — 例", sec["examples"])
        if sec.get("tips"):
            add_bullet_slide(prs, f"{sec['heading']} — 注意点", sec["tips"], eyebrow="Tips")
    if data.get("exercises"):
        add_exercise_slide(prs, data["exercises"])
    if data.get("summary"):
        add_summary_slide(prs, data["summary"])

    out = PPT_DIR / f"p{p:02d}.pptx"
    prs.save(str(out))
    return out


def main() -> None:
    files = sorted(OUTLINES.glob("p*.json"))
    if not files:
        print("no outlines found")
        return
    for f in files:
        out = build_one(f)
        print(f"[ok] {out.name}  ({out.stat().st_size//1024} KB)")
    print(f"BUILD DONE: {len(files)} pptx generated under {PPT_DIR}")


if __name__ == "__main__":
    main()
